import json
import time
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / "results"
FIGURES = RESULTS / "figures"
WEEK_REPORT = ROOT.parents[1] / "reports" / "week5_report.md"


def read_json(name):
    return json.loads((RESULTS / name).read_text(encoding="utf-8"))


def write_text(path, text):
    path.write_text(text, encoding="utf-8")


def set_style():
    plt.rcParams["figure.dpi"] = 140
    plt.rcParams["savefig.dpi"] = 180
    plt.rcParams["font.family"] = "DejaVu Sans"
    plt.style.use("seaborn-v0_8-whitegrid")


def load_metrics():
    all_scales = read_json("week5_all_scales_summary.json")
    features = read_json("week5_features_reproduction_summary.json")
    raw = all_scales["N49000"]
    return {
        "raw_softmax_val": raw["softmax"]["best"]["val_acc"],
        "raw_softmax_test": raw["softmax"]["best"]["test_acc"],
        "raw_twolayer_val": raw["two_layer_net"]["best"]["best_val_acc"],
        "raw_twolayer_test": raw["two_layer_net"]["best"]["test_acc"],
        "feat_softmax_val": features["softmax"]["best"]["val_acc"],
        "feat_softmax_test": features["softmax"]["best"]["test_acc"],
        "feat_twolayer_val": features["two_layer_net"]["best"]["best_val_acc"],
        "feat_twolayer_test": features["two_layer_net"]["best"]["test_acc"],
        "feature_dim": features["feature_info"]["feature_dim"],
        "extract_wall": features["feature_info"]["wall_seconds"],
        "softmax_best": features["softmax"]["best"],
        "twolayer_best": features["two_layer_net"]["best"],
        "twolayer_configs": features["two_layer_net"]["all_configs"],
        "softmax_search": features["softmax"]["search"],
        "total_wall": features["total_wall_seconds"],
    }


def plot_feature_vs_raw(m):
    labels = ["Softmax", "TwoLayerNet"]
    raw = [m["raw_softmax_test"], m["raw_twolayer_test"]]
    feat = [m["feat_softmax_test"], m["feat_twolayer_test"]]
    x = np.arange(len(labels))
    width = 0.34
    fig, ax = plt.subplots(figsize=(8.5, 5))
    bars1 = ax.bar(x - width / 2, raw, width, label="Raw pixels", color="#7b8da8")
    bars2 = ax.bar(x + width / 2, feat, width, label="HOG + color hist", color="#4f9d7a")
    ax.set_title("Raw pixels vs handcrafted image features (N49000)")
    ax.set_ylabel("Test accuracy")
    ax.set_xticks(x, labels)
    ax.set_ylim(0.3, 0.65)
    ax.legend()
    ax.bar_label(bars1, fmt="%.3f", padding=3)
    ax.bar_label(bars2, fmt="%.3f", padding=3)
    plt.tight_layout()
    plt.savefig(FIGURES / "feature_vs_raw_accuracy.png", bbox_inches="tight")
    plt.close()


def plot_feature_softmax_heatmap(m):
    rows = []
    for item in m["softmax_search"].values():
        rows.append({"lr": item["learning_rate"], "reg": item["reg"], "val_acc": item["val_acc"]})
    df = pd.DataFrame(rows)
    pivot = df.pivot(index="reg", columns="lr", values="val_acc").sort_index(ascending=False)
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    im = ax.imshow(pivot.values, cmap="YlGnBu", vmin=pivot.values.min(), vmax=pivot.values.max())
    ax.set_xticks(range(len(pivot.columns)), [f"{x:.1e}" for x in pivot.columns])
    ax.set_yticks(range(len(pivot.index)), [f"{int(x):,}" for x in pivot.index])
    ax.set_xlabel("Learning rate")
    ax.set_ylabel("Regularization")
    ax.set_title("Softmax on HOG + color histogram features")
    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            ax.text(j, i, f"{pivot.values[i, j]:.3f}", ha="center", va="center", fontsize=9)
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    plt.tight_layout()
    plt.savefig(FIGURES / "features_softmax_hyperparam_heatmap.png", bbox_inches="tight")
    plt.close()


def plot_feature_twolayer_configs(m):
    rows = []
    for cfg in m["twolayer_configs"]:
        rows.append(
            {
                "label": f"cfg{cfg['cfg_id']}\nH={cfg['hidden_dim']}\nlr={cfg['learning_rate']:.0e}\nreg={cfg['reg']}",
                "val": cfg["best_val_acc"],
                "test": cfg["test_acc"],
            }
        )
    df = pd.DataFrame(rows)
    x = np.arange(len(df))
    width = 0.38
    fig, ax = plt.subplots(figsize=(10, 5.2))
    b1 = ax.bar(x - width / 2, df["val"], width, label="Validation", color="#4f9d7a")
    b2 = ax.bar(x + width / 2, df["test"], width, label="Test", color="#d78c45")
    ax.set_title("TwoLayerNet tuning on HOG + color histogram features")
    ax.set_ylabel("Accuracy")
    ax.set_xticks(x, df["label"], fontsize=8)
    ax.set_ylim(0.5, 0.64)
    ax.legend()
    ax.bar_label(b1, fmt="%.3f", fontsize=8, padding=2)
    ax.bar_label(b2, fmt="%.3f", fontsize=8, padding=2)
    plt.tight_layout()
    plt.savefig(FIGURES / "features_twolayer_config_comparison.png", bbox_inches="tight")
    plt.close()


def feature_section(m):
    return f"""

## 9. 补充实验：Higher Level Representations - Image Features

Assignment 1 后半部分要求比较 raw pixels 与更高级的手工图像特征。本次已补跑 `features.ipynb` 对应实验：对 CIFAR-10 图像提取 HOG 特征和 HSV color histogram，再在这些特征上训练 Softmax 与 TwoLayerNet。

本次特征设置如下：

- HOG：刻画局部边缘和梯度方向，更接近物体轮廓信息；
- HSV color histogram：刻画图像颜色分布，补充原始像素和边缘特征；
- 特征维度：{m['feature_dim']}；
- 数据规模：49,000 train / 1,000 validation / 1,000 test；
- 特征提取耗时：{m['extract_wall']:.1f}s，总实验耗时：{m['total_wall']:.1f}s。

### 9.1 Raw pixels 与 image features 对比

| 模型 | Raw pixels Test Acc | HOG + Color Hist Test Acc | 提升 |
|---|---:|---:|---:|
| Softmax | {m['raw_softmax_test']:.3f} | {m['feat_softmax_test']:.3f} | +{m['feat_softmax_test'] - m['raw_softmax_test']:.3f} |
| TwoLayerNet | {m['raw_twolayer_test']:.3f} | {m['feat_twolayer_test']:.3f} | +{m['feat_twolayer_test'] - m['raw_twolayer_test']:.3f} |

![Raw pixels 与 image features 对比](figures/feature_vs_raw_accuracy.png)

结果说明：HOG + color histogram 明显优于 raw pixels。Softmax 从 {m['raw_softmax_test']:.3f} 提升到 {m['feat_softmax_test']:.3f}，说明即使模型仍然是线性的，只要输入表示更接近图像语义，分类效果也会提升。TwoLayerNet 从 {m['raw_twolayer_test']:.3f} 提升到 {m['feat_twolayer_test']:.3f}，说明“更好的特征表示 + 非线性分类器”可以叠加带来更强效果。

### 9.2 Softmax on features 调参结果

最佳 Softmax 特征版配置：

```text
learning_rate = {m['softmax_best']['learning_rate']}
reg = {m['softmax_best']['reg']}
val_acc = {m['softmax_best']['val_acc']:.3f}
test_acc = {m['softmax_best']['test_acc']:.3f}
```

![Features Softmax heatmap](figures/features_softmax_hyperparam_heatmap.png)

### 9.3 TwoLayerNet on features 调参结果

最佳 TwoLayerNet 特征版配置：

```text
hidden_dim = {m['twolayer_best']['hidden_dim']}
learning_rate = {m['twolayer_best']['learning_rate']}
reg = {m['twolayer_best']['reg']}
epochs = {m['twolayer_best']['num_epochs']}
val_acc = {m['twolayer_best']['best_val_acc']:.3f}
test_acc = {m['twolayer_best']['test_acc']:.3f}
```

![Features TwoLayerNet config comparison](figures/features_twolayer_config_comparison.png)

这个结果已经达到 CS231n Assignment 1 对 image features 部分常见的目标区间：在 HOG + color histogram 上训练两层网络，测试准确率可以超过 0.58 左右。本地快速测试集上，本次 best-by-validation 模型 test accuracy 为 {m['feat_twolayer_test']:.3f}。

### 9.4 本实验带来的学习点

这部分最重要的结论是：模型性能不只取决于分类器，也强烈取决于输入表示。Raw pixels 把图像当成 3072 维数字向量，很多语义结构被打散；HOG 把边缘和方向编码出来，color histogram 把颜色统计编码出来，因此更接近传统视觉里“可分类”的表示。CNN 后续之所以强，是因为它不再手写 HOG，而是通过卷积层自动学习类似甚至更强的局部视觉特征。
"""


def update_markdown_reports(m):
    section = feature_section(m)
    full_path = RESULTS / "week5_full_reproduction_report.md"
    full_text = full_path.read_text(encoding="utf-8")
    marker = "\n## 9. 补充实验：Higher Level Representations - Image Features"
    if marker in full_text:
        full_text = full_text.split(marker)[0].rstrip() + section
    else:
        full_text = full_text.rstrip() + section
    full_text = full_text.replace(
        "- 仍然使用 raw pixels，没有使用 HOG / color histogram 等特征。",
        "- raw pixels 基线已经完成，并已补充 HOG + color histogram 特征实验。",
    )
    full_text = full_text.replace(
        "1. 跑 `features.ipynb`：加入 HOG 和颜色直方图特征。",
        "1. `features.ipynb` 对应实验已补跑：HOG + color histogram 明显提升 Softmax 和 TwoLayerNet。",
    )
    write_text(full_path, full_text)

    week_text = WEEK_REPORT.read_text(encoding="utf-8")
    week_marker = "\n## 十、补充：Image Features 实验"
    week_section = f"""

## 十、补充：Image Features 实验

在前一版 Week5 报告里，我把 HOG / color histogram 写成后续工作。现在这部分已经补跑完成，对应 CS231n Assignment 1 的 `features.ipynb`。

补充实验使用 49,000 张训练图像、1,000 张验证图像和 1,000 张测试图像。先提取 HOG 与 HSV color histogram，再分别训练 Softmax 和 TwoLayerNet。

| 输入表示 | 模型 | 验证准确率 | 测试准确率 |
|---|---|---:|---:|
| raw pixels | Softmax | {m['raw_softmax_val']:.3f} | {m['raw_softmax_test']:.3f} |
| HOG + color histogram | Softmax | {m['feat_softmax_val']:.3f} | {m['feat_softmax_test']:.3f} |
| raw pixels | TwoLayerNet | {m['raw_twolayer_val']:.3f} | {m['raw_twolayer_test']:.3f} |
| HOG + color histogram | TwoLayerNet | {m['feat_twolayer_val']:.3f} | {m['feat_twolayer_test']:.3f} |

这个结果说明，手工图像特征确实能显著改善分类效果。Softmax 虽然仍是线性模型，但输入换成 HOG + color histogram 后，测试准确率从 {m['raw_softmax_test']:.3f} 提升到 {m['feat_softmax_test']:.3f}。TwoLayerNet 进一步把测试准确率提升到 {m['feat_twolayer_test']:.3f}。这说明 A1 的重点不只是“换模型”，也包括“换表示”：更好的图像表示会让简单分类器也变强。

因此，当前 Week5 的结论需要更新为：本周已经完成 raw pixels 基线、HOG + color histogram 特征实验，以及两类输入表示下的 Softmax / TwoLayerNet 对比。剩余未展开的主要是 FullyConnectedNet 的 dropout / batch normalization、更完整的 10,000 张测试集评估，以及后续 CNN。
"""
    if week_marker in week_text:
        week_text = week_text.split(week_marker)[0].rstrip() + week_section
    else:
        week_text = week_text.rstrip() + week_section
    week_text = week_text.replace(
        "1. 主要使用 raw pixels，没有充分加入 HOG 和 color histogram 等手工特征；",
        "1. raw pixels 基线和 HOG + color histogram 手工特征实验均已完成；",
    )
    write_text(WEEK_REPORT, week_text)


def add_textbox(slide, text, x, y, w, h, size=20, bold=False, color=(42, 46, 52)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_title(slide, kicker, title):
    add_textbox(slide, kicker.upper(), 0.55, 0.28, 2.8, 0.26, size=9, bold=True, color=(95, 111, 130))
    add_textbox(slide, title, 0.55, 0.55, 11.7, 0.55, size=24, bold=True, color=(30, 35, 40))


def add_bullets(slide, bullets, x=0.75, y=1.35, w=5.5, h=4.8, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.space_after = Pt(8)
    return box


def slide_bg(slide, color=(248, 247, 244)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_picture(slide, filename, x, y, w, h=None):
    path = FIGURES / filename
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def update_ppt(m):
    ppt_path = RESULTS / "week5_cs231n_assignment1_report.pptx"
    prs = Presentation(str(ppt_path))
    blank = prs.slide_layouts[6]

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace(
                        "尚未充分跑 HOG / color histogram features",
                        "已补跑 HOG / color histogram features",
                    )
                    run.text = run.text.replace(
                        "1. 跑 features.ipynb：HOG + color histogram",
                        "1. features.ipynb 已补跑：HOG + color histogram",
                    )

    s = prs.slides.add_slide(blank)
    slide_bg(s)
    add_title(s, "image features", "补充实验：HOG + Color Histogram")
    add_bullets(
        s,
        [
            "A1 后半部分要求比较 raw pixels 与 higher-level representations",
            f"本次提取 HOG + HSV color histogram，特征维度 {m['feature_dim']}",
            "HOG 关注边缘和梯度方向，color histogram 关注颜色分布",
            "这一步不是换模型，而是换输入表示",
        ],
        x=0.75,
        y=1.35,
        w=5.1,
        size=15,
    )
    add_picture(s, "feature_vs_raw_accuracy.png", 6.1, 1.2, 6.4, 4.4)

    s = prs.slides.add_slide(blank)
    slide_bg(s)
    add_title(s, "feature tuning", "特征版模型调参结果")
    add_picture(s, "features_softmax_hyperparam_heatmap.png", 0.65, 1.15, 5.8, 4.2)
    add_picture(s, "features_twolayer_config_comparison.png", 6.65, 1.15, 6.0, 4.2)
    add_textbox(
        s,
        f"Best feature TwoLayerNet: val={m['feat_twolayer_val']:.3f}, test={m['feat_twolayer_test']:.3f}",
        0.75,
        6.1,
        11.5,
        0.4,
        size=18,
        bold=True,
        color=(79, 137, 104),
    )

    s = prs.slides.add_slide(blank)
    slide_bg(s)
    add_title(s, "updated conclusion", "更新后的 Week5 结论")
    add_bullets(
        s,
        [
            f"Softmax: raw pixels {m['raw_softmax_test']:.3f} -> features {m['feat_softmax_test']:.3f}",
            f"TwoLayerNet: raw pixels {m['raw_twolayer_test']:.3f} -> features {m['feat_twolayer_test']:.3f}",
            "手工图像特征能显著提升传统分类器，因为它更接近边缘、颜色、轮廓等视觉语义",
            "CNN 后续要学习的事情，本质上就是自动学习比 HOG 更强的局部图像特征",
        ],
        x=0.8,
        y=1.35,
        w=11.4,
        size=18,
    )

    prs.save(ppt_path)


def update_speaker_notes(m):
    notes_path = RESULTS / "week5_cs231n_assignment1_speaker_notes.md"
    text = notes_path.read_text(encoding="utf-8")
    marker = "\n## Slide 15"
    addition = f"""

## Slide 15：补充实验 - HOG + Color Histogram

这一页要说明，Assignment 1 后半部分不只是用 raw pixels，还要求我们尝试 higher-level representations。我补跑了 HOG 和 HSV color histogram。HOG 主要描述边缘和梯度方向，color histogram 描述颜色分布。它们不是深度学习自动学出来的特征，而是传统计算机视觉里手工设计的图像表示。

这里最关键的对比是：Softmax 在 raw pixels 上测试准确率是 {m['raw_softmax_test']:.3f}，换成 HOG + color histogram 后提升到 {m['feat_softmax_test']:.3f}；TwoLayerNet 在 raw pixels 上是 {m['raw_twolayer_test']:.3f}，换成特征后提升到 {m['feat_twolayer_test']:.3f}。这说明输入表示本身非常重要。

## Slide 16：特征版调参结果

这一页展示两个调参图。左边是 Softmax 在特征上的 learning rate 和 regularization 搜索，右边是 TwoLayerNet 不同 hidden_dim、learning_rate 和 regularization 配置的比较。最佳 TwoLayerNet 特征版模型验证准确率达到 {m['feat_twolayer_val']:.3f}，测试准确率达到 {m['feat_twolayer_test']:.3f}。

这里可以强调，我们选择模型仍然依据 validation accuracy，而不是直接挑 test accuracy 最高的配置。这和前面 raw pixels 实验保持一致。

## Slide 17：更新后的结论

补充 features 实验之后，Week5 的结论更完整了。第一，raw pixels 可以跑通完整 pipeline，但表达能力有限。第二，HOG + color histogram 这样的手工特征能显著提升 Softmax 和 TwoLayerNet。第三，后续 CNN 的意义就更清楚了：CNN 不再依赖人手写 HOG，而是通过卷积层自动学习局部边缘、纹理和更高级的视觉模式。
"""
    if marker in text:
        text = text.split(marker)[0].rstrip() + addition
    else:
        text = text.rstrip() + addition
    write_text(notes_path, text)


def main():
    set_style()
    m = load_metrics()
    plot_feature_vs_raw(m)
    plot_feature_softmax_heatmap(m)
    plot_feature_twolayer_configs(m)
    update_markdown_reports(m)
    update_ppt(m)
    update_speaker_notes(m)
    print("updated feature report assets")
    print(time.strftime("%Y-%m-%d %H:%M:%S"))


if __name__ == "__main__":
    main()
