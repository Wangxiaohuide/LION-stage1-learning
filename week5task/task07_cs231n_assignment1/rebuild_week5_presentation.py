import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / "results"
FIGURES = RESULTS / "figures"
OUT = RESULTS / "week5_cs231n_assignment1_report.pptx"
NOTES = RESULTS / "week5_cs231n_assignment1_speaker_notes.md"
OUTLINE = RESULTS / "week5_cs231n_assignment1_ppt_outline.md"


def read_json(name):
    return json.loads((RESULTS / name).read_text(encoding="utf-8"))


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=(35, 40, 46), align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    if align:
        p.alignment = align
    return box


def add_title(slide, section, title):
    add_text(slide, section.upper(), 0.55, 0.25, 3.3, 0.28, 9, True, (92, 107, 126))
    add_text(slide, title, 0.55, 0.58, 11.9, 0.55, 24, True, (26, 31, 36))


def add_bullets(slide, bullets, x, y, w, h, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(42, 46, 52)
        p.space_after = Pt(8)
    return box


def add_image(slide, filename, x, y, w, h=None):
    path = FIGURES / filename
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_metric(slide, label, value, x, y, color=(79, 137, 104)):
    add_text(slide, label, x, y, 2.7, 0.25, 11, True, (92, 107, 126), PP_ALIGN.CENTER)
    add_text(slide, value, x, y + 0.28, 2.7, 0.5, 24, True, color, PP_ALIGN.CENTER)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 247, 244)


def make_table(slide, rows, x, y, w, h, font_size=11):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [w / len(rows[0])] * len(rows[0])
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(235, 239, 236) if r == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = RGBColor(35, 40, 46)
    return table


def build():
    all_results = read_json("week5_all_scales_summary.json")
    feat = read_json("week5_features_reproduction_summary.json")
    n49 = all_results["N49000"]

    raw_softmax = n49["softmax"]["best"]
    raw_tn = n49["two_layer_net"]["best"]
    feat_softmax = feat["softmax"]["best"]
    feat_tn = feat["two_layer_net"]["best"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank); bg(s)
    add_text(s, "Week5 CS231n Assignment 1", 0.7, 0.95, 10.8, 0.6, 34, True)
    add_text(s, "CIFAR-10 图像分类本地复现：raw pixels 基线 + HOG/color histogram 特征实验", 0.72, 1.72, 11.4, 0.48, 18)
    add_metric(s, "最佳 raw pixels", f"{raw_tn['test_acc']:.3f}", 0.75, 2.65)
    add_metric(s, "最佳 image features", f"{feat_tn['test_acc']:.3f}", 3.65, 2.65)
    add_metric(s, "features 提升", f"+{feat_tn['test_acc'] - raw_tn['test_acc']:.3f}", 6.55, 2.65, (205, 117, 57))
    add_bullets(s, [
        "目标：在本地复现 CS231n A1 的完整学习流程，而不是提交官方作业",
        "主线：kNN、Softmax、TwoLayerNet、HOG + color histogram",
        "汇报重点：为什么“输入表示”会显著影响图像分类效果",
    ], 0.85, 4.0, 11.4, 1.7, 17)

    # 2
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "assignment 1 scope", "Assignment 1 到底要求我们理解什么")
    make_table(s, [
        ["题目", "内容", "本次是否完成"],
        ["Q1", "kNN classifier", "已完成"],
        ["Q2", "Softmax classifier", "已完成"],
        ["Q3", "Two-Layer Neural Network", "已完成"],
        ["Q4", "Higher-level image features", "已补跑"],
        ["Q5", "FullyConnectedNet", "基础验证已跑，dropout/BN 未展开"],
    ], 0.8, 1.35, 11.7, 4.4, 12)
    add_text(s, "这次新增的就是 Q4：不用 raw pixels，而是加入 HOG 和 color histogram。", 0.9, 6.15, 11.3, 0.45, 17, True, (79, 137, 104))

    # 3
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "dataset", "CIFAR-10 与实验划分")
    add_bullets(s, [
        "CIFAR-10：10 类 32×32 彩色图像",
        "官方数据：50,000 train + 10,000 test",
        "A1 常用划分：49,000 train + 1,000 validation",
        "本地快速评估：使用 1,000 张 test 做模型对比",
    ], 0.75, 1.35, 5.3, 4.4, 17)
    add_image(s, "data_scale_effect.png", 6.35, 1.25, 6.0, 4.2)

    # 4
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "workflow", "实验流程：训练、调参、验证、测试、可视化")
    add_bullets(s, [
        "1. 下载 CIFAR-10 并构造 train / validation / test split",
        "2. 先跑 raw pixels 基线：kNN、Softmax、TwoLayerNet",
        "3. 再跑 image features：HOG + HSV color histogram",
        "4. 所有模型用 validation accuracy 选超参数",
        "5. 最后只用 test accuracy 汇报泛化表现",
    ], 0.9, 1.35, 11.2, 4.8, 19)

    # 5
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "raw pixels", "Raw pixels 基线结果")
    rows = [
        ["模型", "最佳参数", "Val Acc", "Test Acc"],
        ["kNN", f"k={n49['knn']['best_k']}", f"{n49['knn']['best_val_acc']:.3f}", f"{n49['knn']['test_acc']:.3f}"],
        ["Softmax", f"lr={raw_softmax['learning_rate']:.1e}, reg={raw_softmax['reg']:.0f}", f"{raw_softmax['val_acc']:.3f}", f"{raw_softmax['test_acc']:.3f}"],
        ["TwoLayerNet", f"H={raw_tn['hidden_dim']}, lr={raw_tn['learning_rate']:.1e}", f"{raw_tn['best_val_acc']:.3f}", f"{raw_tn['test_acc']:.3f}"],
    ]
    make_table(s, rows, 0.75, 1.35, 11.8, 2.1, 12)
    add_image(s, "model_test_accuracy_comparison.png", 1.25, 3.7, 10.7, 3.0)

    # 6
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "raw pixels limits", "为什么 raw pixels 表现有限")
    add_bullets(s, [
        "32×32×3 图像被拉平成 3072 维向量后，空间结构被打散",
        "kNN 的 L2 距离容易受背景、亮度、位置变化影响",
        "Softmax 只能学习线性类别模板，决策边界表达能力有限",
        "TwoLayerNet 加入非线性后提升明显，但仍没有显式利用局部图像结构",
    ], 0.85, 1.35, 11.4, 4.5, 20)

    # 7
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "image features", "补充实验：HOG + Color Histogram")
    add_bullets(s, [
        "HOG：统计局部梯度方向，突出边缘与轮廓",
        "HSV color histogram：统计颜色分布，补充颜色语义",
        f"提取后特征维度为 {feat['feature_info']['feature_dim']}，远小于 raw pixels 的 3072 维",
        "这一步的核心不是换分类器，而是换输入表示",
    ], 0.75, 1.35, 5.4, 4.4, 17)
    add_image(s, "feature_vs_raw_accuracy.png", 6.35, 1.25, 6.0, 4.25)

    # 8
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "feature results", "Raw pixels vs Image Features")
    make_table(s, [
        ["输入表示", "模型", "Val Acc", "Test Acc"],
        ["raw pixels", "Softmax", f"{raw_softmax['val_acc']:.3f}", f"{raw_softmax['test_acc']:.3f}"],
        ["HOG + color hist", "Softmax", f"{feat_softmax['val_acc']:.3f}", f"{feat_softmax['test_acc']:.3f}"],
        ["raw pixels", "TwoLayerNet", f"{raw_tn['best_val_acc']:.3f}", f"{raw_tn['test_acc']:.3f}"],
        ["HOG + color hist", "TwoLayerNet", f"{feat_tn['best_val_acc']:.3f}", f"{feat_tn['test_acc']:.3f}"],
    ], 0.8, 1.3, 11.7, 2.7, 12)
    add_text(s, f"结论：TwoLayerNet + image features 的测试准确率达到 {feat_tn['test_acc']:.3f}，比 raw pixels 提升 {feat_tn['test_acc'] - raw_tn['test_acc']:.3f}。", 0.9, 4.6, 11.2, 0.55, 19, True, (79, 137, 104))

    # 9
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "softmax tuning", "Softmax on Features：学习率与正则化搜索")
    add_image(s, "features_softmax_hyperparam_heatmap.png", 0.85, 1.2, 7.0, 4.8)
    add_bullets(s, [
        f"最佳 lr={feat_softmax['learning_rate']:.1e}",
        f"最佳 reg={feat_softmax['reg']:.0f}",
        f"Val Acc={feat_softmax['val_acc']:.3f}",
        f"Test Acc={feat_softmax['test_acc']:.3f}",
        "线性模型也能受益于更好的图像表示",
    ], 8.25, 1.55, 4.1, 4.1, 17)

    # 10
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "two-layer tuning", "TwoLayerNet on Features：非线性分类器继续提升")
    add_image(s, "features_twolayer_config_comparison.png", 0.7, 1.2, 7.3, 4.5)
    add_bullets(s, [
        f"最佳 hidden_dim={feat_tn['hidden_dim']}",
        f"最佳 lr={feat_tn['learning_rate']:.1e}",
        f"最佳 reg={feat_tn['reg']}",
        f"Val Acc={feat_tn['best_val_acc']:.3f}",
        f"Test Acc={feat_tn['test_acc']:.3f}",
    ], 8.35, 1.55, 4.0, 4.0, 17)

    # 11
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "interpretation", "这个补充实验说明了什么")
    add_bullets(s, [
        "图像分类性能不仅取决于模型，也取决于输入表示",
        "HOG 把边缘和轮廓编码出来，让分类器看到更接近语义的信号",
        "颜色直方图保留类别相关的颜色统计，比如 frog / ship / truck 的色彩分布差异",
        "CNN 的意义更清楚了：它会自动学习比 HOG 更强、更层次化的局部视觉特征",
    ], 0.85, 1.35, 11.4, 4.7, 20)

    # 12
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "takeaways", "本周学习收获")
    add_bullets(s, [
        "掌握了 CIFAR-10 图像分类 pipeline：数据划分、训练、验证调参、测试评估",
        "理解了 kNN、Softmax、TwoLayerNet 的表达能力差异",
        "练习了超参数搜索：learning rate、regularization、hidden_dim、epochs",
        "通过 HOG + color histogram 理解“特征表示”对模型效果的影响",
        "为后续 CNN 学习建立了动机：从手工特征走向自动特征学习",
    ], 0.85, 1.35, 11.4, 4.8, 19)

    # 13
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "current limits", "当前仍然保留的局限")
    add_bullets(s, [
        "测试集仍使用 1,000 张快速评估，后续可扩到完整 10,000 张",
        "FullyConnectedNet 的 dropout / batch normalization 还没有系统展开",
        "TwoLayerNet 调参还可以继续更细，例如更多 epoch 和更密集搜索",
        "目前还没有训练 CNN，因此还没有自动利用图像局部空间结构",
    ], 0.85, 1.35, 11.4, 4.8, 20)

    # 14
    s = prs.slides.add_slide(blank); bg(s); add_title(s, "next step", "下一步：从 Assignment 1 走向 Assignment 2")
    add_bullets(s, [
        "继续 FullyConnectedNet：多层网络、Dropout、Batch Normalization",
        "进入 CNN：卷积层、池化层、局部感受野、权重共享",
        "用 PyTorch 重新实现 CIFAR-10 训练流程",
        "目标从“手工特征 + 小网络”过渡到“端到端自动学习图像特征”",
    ], 0.85, 1.35, 11.4, 4.8, 20)

    # 15
    s = prs.slides.add_slide(blank); bg(s)
    add_text(s, "汇报结论", 0.75, 1.1, 10.5, 0.65, 34, True)
    add_text(s, "Assignment 1 的价值不是只跑出一个准确率，而是理解图像分类系统由数据、表示、模型、优化和调参共同决定。", 0.8, 2.1, 11.6, 0.9, 23, True, (35, 40, 46))
    add_bullets(s, [
        f"raw pixels 最佳 TwoLayerNet test acc：{raw_tn['test_acc']:.3f}",
        f"HOG + color histogram 最佳 TwoLayerNet test acc：{feat_tn['test_acc']:.3f}",
        "补充 features 后，Week5 已覆盖 A1 的主要实验主线",
    ], 0.95, 3.45, 11.2, 2.0, 21)

    prs.save(OUT)
    write_notes(raw_softmax, raw_tn, feat_softmax, feat_tn, feat)
    write_outline()


def write_notes(raw_softmax, raw_tn, feat_softmax, feat_tn, feat):
    notes = f"""# Week5 CS231n Assignment 1 PPT 逐页讲稿

## Slide 1：标题页
今天汇报 Week5 的 CS231n Assignment 1 本地复现实验。重点不是提交官方作业，而是在本地复现 CIFAR-10 图像分类流程。最终结果是：raw pixels 上最佳 TwoLayerNet 测试准确率为 {raw_tn['test_acc']:.3f}；加入 HOG + color histogram 后，最佳 TwoLayerNet 测试准确率提升到 {feat_tn['test_acc']:.3f}。

## Slide 2：Assignment 1 范围
Assignment 1 分成五块：kNN、Softmax、TwoLayerNet、Image Features 和 FullyConnectedNet。本次已经补齐 Q4，也就是 higher-level image features。Q5 的基础验证已经跑过，但 dropout 和 batch normalization 更适合放到 Assignment 2 继续展开。

## Slide 3：数据集与划分
实验使用 CIFAR-10，每张图片是 32×32×3。官方训练集 50,000 张，本次按 A1 常用方式划出 49,000 张训练和 1,000 张验证。测试为了本地快速评估，使用 1,000 张。

## Slide 4：实验流程
整个 pipeline 是先准备数据，再训练模型，再用验证集选择超参数，最后用测试集报告泛化性能。这里强调不能用测试集调参，否则测试结果就不再客观。

## Slide 5：Raw pixels 基线
raw pixels 基线中，kNN 测试准确率 0.348，Softmax 是 {raw_softmax['test_acc']:.3f}，TwoLayerNet 是 {raw_tn['test_acc']:.3f}。结果说明非线性隐藏层确实比线性模型更强。

## Slide 6：Raw pixels 的局限
raw pixels 的问题是图像空间结构被打散。一个小的平移、背景变化或者亮度变化，都会让像素距离发生很大变化。所以 kNN 和 Softmax 在 raw pixels 上都有明显上限。

## Slide 7：Image Features 补充实验
这一页是这次新增重点。HOG 统计边缘和梯度方向，color histogram 统计颜色分布。它们让模型看到的不是零散像素，而是更接近视觉语义的手工特征。

## Slide 8：Raw pixels vs Image Features
表格显示，Softmax 从 {raw_softmax['test_acc']:.3f} 提升到 {feat_softmax['test_acc']:.3f}，TwoLayerNet 从 {raw_tn['test_acc']:.3f} 提升到 {feat_tn['test_acc']:.3f}。这说明表示学习或特征工程对图像分类非常重要。

## Slide 9：Softmax on Features
左边热力图展示 Softmax 在特征上的调参。最佳配置是 lr={feat_softmax['learning_rate']:.1e}, reg={feat_softmax['reg']:.0f}，测试准确率 {feat_softmax['test_acc']:.3f}。线性模型本身没变，但输入表示变好后结果明显提升。

## Slide 10：TwoLayerNet on Features
TwoLayerNet 在特征上继续提升。最佳配置是 hidden_dim={feat_tn['hidden_dim']}, lr={feat_tn['learning_rate']:.1e}, reg={feat_tn['reg']}，验证准确率 {feat_tn['best_val_acc']:.3f}，测试准确率 {feat_tn['test_acc']:.3f}。

## Slide 11：实验解释
这一页解释为什么 features 有用。HOG 把边缘和轮廓编码出来，颜色直方图保留颜色统计。CNN 后续要做的，就是不再手工写这些特征，而是从数据中自动学习更强的局部视觉特征。

## Slide 12：学习收获
这周主要学习了图像分类 pipeline、模型表达能力、超参数搜索、验证集和测试集的区别，以及输入表示对模型性能的影响。

## Slide 13：当前局限
现在的局限不再是没有跑 HOG/color histogram，因为这部分已经补上了。剩余局限包括：测试集还是 1,000 张，FullyConnectedNet 的 dropout 和 batch normalization 没有系统展开，还没有进入 CNN。

## Slide 14：下一步
下一步应该进入 Assignment 2：Dropout、BatchNorm、CNN 和 PyTorch。路线是从手工特征过渡到端到端自动学习图像特征。

## Slide 15：汇报结论
最后总结：Assignment 1 的价值不是只得到一个准确率，而是理解图像分类系统由数据、表示、模型、优化和调参共同决定。补充 features 后，Week5 已经覆盖 A1 的主要实验主线。
"""
    NOTES.write_text(notes, encoding="utf-8")


def write_outline():
    outline = """# Week5 CS231n Assignment 1 PPT 大纲

1. 标题页：raw pixels 基线 + image features 补充实验
2. Assignment 1 范围：Q1-Q5 完成情况
3. CIFAR-10 数据集与划分
4. 实验流程
5. Raw pixels 基线结果
6. Raw pixels 的局限
7. Image Features：HOG + Color Histogram
8. Raw pixels vs Image Features 对比
9. Softmax on Features 调参
10. TwoLayerNet on Features 调参
11. 为什么 features 有用
12. 本周学习收获
13. 当前局限
14. 下一步：走向 Assignment 2
15. 汇报结论
"""
    OUTLINE.write_text(outline, encoding="utf-8")


if __name__ == "__main__":
    build()
    print(OUT)
    print(NOTES)
