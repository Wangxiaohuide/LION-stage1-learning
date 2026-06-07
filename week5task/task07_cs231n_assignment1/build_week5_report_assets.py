import json
import time
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.font_manager as font_manager
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from cs231n.classifiers.fc_net import TwoLayerNet
from cs231n.data_utils import get_CIFAR10_data
from cs231n.solver import Solver


ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / "results"
FIGURES = RESULTS / "figures"
FIGURES.mkdir(parents=True, exist_ok=True)

CLASS_NAMES = ["plane", "car", "bird", "cat", "deer", "dog", "frog", "horse", "ship", "truck"]
SCALE_ORDER = ["N10000", "N20000", "N49000"]
MODEL_ORDER = ["kNN", "Softmax", "TwoLayerNet"]


def read_json(path):
    return json.loads(Path(path).read_text(encoding="utf-8"))


def write_json(path, obj):
    Path(path).write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")


def first_existing(*names):
    for name in names:
        p = RESULTS / name
        if p.exists():
            return p
    return None


def normalize_softmax_batch(raw):
    best = raw["batch_best"]
    if "learning_rate" not in best:
        best["learning_rate"] = best.get("lr")
    return raw


def load_all_results():
    full = read_json(RESULTS / "week5_full_reproduction_summary.json")
    all_results = {
        "N10000": full["N10000"],
        "N20000": full["N20000"],
    }

    knn49 = read_json(RESULTS / "knn_49k.json")
    soft_batches = [
        normalize_softmax_batch(read_json(RESULTS / f"softmax_49k_b{i}.json"))
        for i in range(1, 5)
    ]
    soft_search = {}
    soft_best = None
    soft_wall = 0.0
    for batch in soft_batches:
        soft_search.update(batch["search"])
        soft_wall += batch.get("wall_seconds", 0.0)
        b = batch["batch_best"]
        if soft_best is None or b["val_acc"] > soft_best["val_acc"]:
            soft_best = {
                "learning_rate": b["learning_rate"],
                "reg": b["reg"],
                "train_acc": b["train_acc"],
                "val_acc": b["val_acc"],
                "test_acc": b["test_acc"],
            }

    tn_cfgs = [read_json(RESULTS / f"tn49_cfg{i}.json") for i in range(1, 9)]
    normalized_cfgs = []
    for cfg in tn_cfgs:
        normalized_cfgs.append(
            {
                "cfg_id": cfg["cfg_id"],
                "hidden_dim": cfg["hidden_dim"],
                "learning_rate": cfg.get("learning_rate", cfg.get("lr")),
                "reg": cfg["reg"],
                "num_epochs": cfg.get("num_epochs", cfg.get("epochs")),
                "batch_size": cfg["batch_size"],
                "best_val_acc": cfg["best_val_acc"],
                "last_train_acc": cfg["last_train_acc"],
                "test_acc": cfg["test_acc"],
                "wall_seconds": cfg.get("wall_seconds", cfg.get("wall")),
            }
        )
    tn_best = max(normalized_cfgs, key=lambda x: x["best_val_acc"])
    all_results["N49000"] = {
        "num_training": 49000,
        "num_validation": 1000,
        "num_test": 1000,
        "knn": knn49,
        "softmax": {
            "scale": "N49000",
            "num_training": 49000,
            "num_val": 1000,
            "num_test": 1000,
            "search": soft_search,
            "best": soft_best,
            "wall_seconds": round(soft_wall, 2),
        },
        "two_layer_net": {
            "scale": "N49000",
            "num_training": 49000,
            "num_val": 1000,
            "num_test": 1000,
            "all_configs": normalized_cfgs,
            "best": tn_best,
            "wall_seconds": round(sum(c["wall_seconds"] for c in normalized_cfgs), 2),
        },
    }
    return all_results


def metric_rows(all_results):
    rows = []
    for scale in SCALE_ORDER:
        r = all_results[scale]
        rows.append(
            {
                "scale": scale,
                "model": "kNN",
                "val_acc": r["knn"]["best_val_acc"],
                "test_acc": r["knn"]["test_acc"],
                "runtime": r["knn"].get("wall_seconds", 0),
                "params": f"k={r['knn']['best_k']}",
            }
        )
        sm = r["softmax"]["best"]
        rows.append(
            {
                "scale": scale,
                "model": "Softmax",
                "val_acc": sm["val_acc"],
                "test_acc": sm["test_acc"],
                "runtime": r["softmax"].get("wall_seconds", 0),
                "params": f"lr={sm['learning_rate']:.1e}, reg={sm['reg']:.0f}",
            }
        )
        tn = r["two_layer_net"]["best"]
        val_key = "best_val_acc" if "best_val_acc" in tn else "val_acc"
        rows.append(
            {
                "scale": scale,
                "model": "TwoLayerNet",
                "val_acc": tn[val_key],
                "test_acc": tn["test_acc"],
                "runtime": r["two_layer_net"].get("wall_seconds", 0),
                "params": f"H={tn['hidden_dim']}, lr={tn['learning_rate']:.1e}, reg={tn['reg']}",
            }
        )
    return pd.DataFrame(rows)


def set_plot_style():
    for font_path in [
        r"C:\Windows\Fonts\NotoSansSC-VF.ttf",
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
    ]:
        if Path(font_path).exists():
            font_manager.fontManager.addfont(font_path)
    plt.rcParams["figure.dpi"] = 140
    plt.rcParams["savefig.dpi"] = 180
    plt.rcParams["font.family"] = "sans-serif"
    plt.rcParams["font.sans-serif"] = ["Noto Sans SC", "Microsoft YaHei", "SimHei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    plt.style.use("seaborn-v0_8-whitegrid")


def grouped_bar(df, metric, title, ylabel, path):
    pivot = df.pivot(index="scale", columns="model", values=metric).loc[SCALE_ORDER, MODEL_ORDER]
    ax = pivot.plot(kind="bar", figsize=(9, 5), width=0.78, color=["#7b8da8", "#d78c45", "#4f9d7a"])
    ax.set_title(title, fontsize=15, pad=12)
    ax.set_xlabel("Training scale")
    ax.set_ylabel(ylabel)
    ax.set_ylim(0, max(0.6, pivot.to_numpy().max() + 0.08))
    ax.legend(title="")
    for container in ax.containers:
        ax.bar_label(container, fmt="%.3f", fontsize=8, padding=2)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close()


def data_scale_line(df, path):
    fig, ax = plt.subplots(figsize=(9, 5))
    scale_x = [10000, 20000, 49000]
    for model, color in zip(MODEL_ORDER, ["#7b8da8", "#d78c45", "#4f9d7a"]):
        vals = [df[(df.scale == s) & (df.model == model)]["test_acc"].iloc[0] for s in SCALE_ORDER]
        ax.plot(scale_x, vals, marker="o", linewidth=2.4, label=model, color=color)
        for x, y in zip(scale_x, vals):
            ax.text(x, y + 0.012, f"{y:.3f}", ha="center", fontsize=9)
    ax.set_title("Effect of training scale on test accuracy", fontsize=15, pad=12)
    ax.set_xlabel("Number of training samples")
    ax.set_ylabel("Test accuracy")
    ax.set_ylim(0.25, 0.58)
    ax.legend()
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close()


def runtime_chart(df, path):
    pivot = df.pivot(index="scale", columns="model", values="runtime").loc[SCALE_ORDER, MODEL_ORDER]
    ax = pivot.plot(kind="bar", figsize=(9, 5), width=0.78, color=["#7b8da8", "#d78c45", "#4f9d7a"])
    ax.set_title("Runtime comparison", fontsize=15, pad=12)
    ax.set_xlabel("Training scale")
    ax.set_ylabel("Wall seconds")
    ax.legend(title="")
    for container in ax.containers:
        ax.bar_label(container, fmt="%.0f", fontsize=8, padding=2)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close()


def softmax_heatmap(all_results, path):
    search = all_results["N49000"]["softmax"]["search"]
    rows = []
    for item in search.values():
        rows.append({"lr": item["lr"], "reg": item["reg"], "val_acc": item["val_acc"]})
    df = pd.DataFrame(rows)
    pivot = df.pivot(index="reg", columns="lr", values="val_acc").sort_index(ascending=False)
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    im = ax.imshow(pivot.values, cmap="YlGnBu", vmin=pivot.values.min(), vmax=pivot.values.max())
    ax.set_xticks(range(len(pivot.columns)), [f"{x:.1e}" for x in pivot.columns])
    ax.set_yticks(range(len(pivot.index)), [f"{int(x):,}" for x in pivot.index])
    ax.set_xlabel("Learning rate")
    ax.set_ylabel("Regularization")
    ax.set_title("Softmax hyperparameter heatmap (N49000, val accuracy)", fontsize=14, pad=12)
    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            ax.text(j, i, f"{pivot.values[i, j]:.3f}", ha="center", va="center", fontsize=9)
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close()


def two_layer_config_chart(all_results, path):
    rows = []
    for cfg in all_results["N49000"]["two_layer_net"]["all_configs"]:
        label = f"cfg{cfg['cfg_id']}\nH={cfg['hidden_dim']}\nlr={cfg['learning_rate']:.1e}\nreg={cfg['reg']}"
        rows.append({"label": label, "val": cfg["best_val_acc"], "test": cfg["test_acc"]})
    df = pd.DataFrame(rows)
    x = np.arange(len(df))
    width = 0.38
    fig, ax = plt.subplots(figsize=(11, 5.4))
    ax.bar(x - width / 2, df["val"], width, label="Validation", color="#6b8fb9")
    ax.bar(x + width / 2, df["test"], width, label="Test", color="#4f9d7a")
    ax.set_title("TwoLayerNet configuration comparison (N49000)", fontsize=15, pad=12)
    ax.set_ylabel("Accuracy")
    ax.set_xticks(x, df["label"], fontsize=8)
    ax.set_ylim(0.40, 0.55)
    ax.legend()
    for container in ax.containers:
        ax.bar_label(container, fmt="%.3f", fontsize=8, padding=2)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close()


def train_best_for_diagnostics():
    diag_json = RESULTS / "best_twolayer_49k_diagnostics.json"
    if diag_json.exists() and (FIGURES / "two_layer_loss_curve.png").exists() and (FIGURES / "confusion_matrix_best_model.png").exists():
        return read_json(diag_json)

    np.random.seed(231)
    data = get_CIFAR10_data(num_training=49000, num_validation=1000, num_test=1000)
    cfg = {"hidden_dim": 100, "learning_rate": 2.5e-4, "reg": 0.05, "num_epochs": 8, "batch_size": 512}
    t0 = time.time()
    model = TwoLayerNet(input_dim=3 * 32 * 32, hidden_dim=cfg["hidden_dim"], num_classes=10, weight_scale=1e-2, reg=cfg["reg"])
    solver = Solver(
        model,
        data,
        update_rule="adam",
        optim_config={"learning_rate": cfg["learning_rate"]},
        lr_decay=0.95,
        num_epochs=cfg["num_epochs"],
        batch_size=cfg["batch_size"],
        print_every=9999,
        verbose=False,
        num_train_samples=1000,
        num_val_samples=None,
    )
    solver.train()
    scores = model.loss(data["X_test"])
    pred = np.argmax(scores, axis=1)
    y = data["y_test"]
    cm = np.zeros((10, 10), dtype=int)
    for truth, guess in zip(y, pred):
        cm[int(truth), int(guess)] += 1

    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.plot(solver.loss_history, color="#4f9d7a", linewidth=2)
    ax.set_title("Best TwoLayerNet loss curve (N49000)", fontsize=14, pad=12)
    ax.set_xlabel("Iteration")
    ax.set_ylabel("Loss")
    plt.tight_layout()
    plt.savefig(FIGURES / "two_layer_loss_curve.png", bbox_inches="tight")
    plt.close()

    fig, ax = plt.subplots(figsize=(7.2, 6.2))
    im = ax.imshow(cm, cmap="Blues")
    ax.set_title("Best TwoLayerNet confusion matrix (1000 test images)", fontsize=14, pad=12)
    ax.set_xlabel("Predicted")
    ax.set_ylabel("True")
    ax.set_xticks(range(10), CLASS_NAMES, rotation=45, ha="right")
    ax.set_yticks(range(10), CLASS_NAMES)
    for i in range(10):
        for j in range(10):
            ax.text(j, i, str(cm[i, j]), ha="center", va="center", fontsize=7, color="white" if cm[i, j] > cm.max() * 0.55 else "#1d2a35")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    plt.tight_layout()
    plt.savefig(FIGURES / "confusion_matrix_best_model.png", bbox_inches="tight")
    plt.close()

    diag = {
        "config": cfg,
        "best_val_acc": float(solver.best_val_acc),
        "test_acc": float(np.mean(pred == y)),
        "wall_seconds": round(time.time() - t0, 2),
        "loss_history_len": len(solver.loss_history),
    }
    write_json(diag_json, diag)
    return diag


def make_figures(all_results):
    set_plot_style()
    df = metric_rows(all_results)
    grouped_bar(df, "test_acc", "Model test accuracy comparison", "Test accuracy", FIGURES / "model_test_accuracy_comparison.png")
    grouped_bar(df, "val_acc", "Model validation accuracy comparison", "Validation accuracy", FIGURES / "model_val_accuracy_comparison.png")
    softmax_heatmap(all_results, FIGURES / "softmax_hyperparam_heatmap.png")
    two_layer_config_chart(all_results, FIGURES / "two_layer_config_comparison.png")
    data_scale_line(df, FIGURES / "data_scale_effect.png")
    runtime_chart(df, FIGURES / "runtime_comparison.png")
    diag = train_best_for_diagnostics()
    return df, diag


def best_overall(df):
    return df.sort_values("test_acc", ascending=False).iloc[0].to_dict()


def write_report(all_results, df, diag):
    best = best_overall(df)
    soft_best = all_results["N49000"]["softmax"]["best"]
    tn_best = all_results["N49000"]["two_layer_net"]["best"]
    md = f"""# Week5 CS231n Assignment 1 本地复现实验报告

生成时间：{time.strftime("%Y-%m-%d %H:%M:%S")}

## 1. 实验目标

本实验在本地复现 CS231n Assignment 1 的核心图像分类流程。目标不是提交作业，而是用 CIFAR-10 数据集理解从最近邻、线性分类器到两层神经网络的完整 pipeline：数据加载、训练/验证/测试划分、超参数搜索、模型评估和结果可视化。

## 2. 数据集与规模

CIFAR-10 包含 10 个类别、32×32 彩色图像。官方常用作业划分为：

- 训练集：49,000 张
- 验证集：1,000 张
- 测试集：本次快速复现实验使用 1,000 张

本报告对比了三个训练规模：N10000、N20000、N49000。

## 3. 模型与调参设置

- kNN：基于原始像素空间 L2 距离，调节 k 值。
- Softmax：原始像素展平后训练线性分类器，搜索 learning rate 与 regularization strength。
- TwoLayerNet：一层隐藏层 + ReLU + Softmax，使用 Adam 优化器，搜索 hidden_dim、learning_rate、reg 和 epoch 数。

## 4. 核心结果

| 规模 | 模型 | 最佳参数 | Val Acc | Test Acc | 耗时(s) |
|---|---|---|---:|---:|---:|
"""
    for _, row in df.iterrows():
        md += f"| {row['scale']} | {row['model']} | {row['params']} | {row['val_acc']:.3f} | {row['test_acc']:.3f} | {row['runtime']:.1f} |\n"

    md += f"""

整体最佳模型按测试集准确率排序为：**{best['scale']} / {best['model']}**，测试准确率 **{best['test_acc']:.3f}**。

需要注意：按验证集选择模型时，N49000 的最佳 TwoLayerNet 是 cfg4：hidden=100、lr=2.5e-4、reg=0.05、8 epochs，validation accuracy 为 {tn_best['best_val_acc']:.3f}，test accuracy 为 {tn_best['test_acc']:.3f}。另一个 cfg5 在测试集上达到 0.510，但验证集略低，说明测试集上的单次波动不能替代验证集模型选择。

## 5. 可视化结果

### 5.1 模型测试准确率对比

![模型测试准确率对比](figures/model_test_accuracy_comparison.png)

TwoLayerNet 明显优于 Softmax 和 kNN。kNN 在原始像素空间中受限最大，Softmax 能学到线性类别模板，而 TwoLayerNet 引入非线性隐藏层后表达能力更强。

### 5.2 模型验证集准确率对比

![模型验证集准确率对比](figures/model_val_accuracy_comparison.png)

验证集结果更适合做模型选择。N49000 下 TwoLayerNet 的最佳验证准确率超过 0.51，说明更多训练数据和适当正则化能明显改善泛化。

### 5.3 Softmax 超参数热力图

![Softmax 超参数热力图](figures/softmax_hyperparam_heatmap.png)

N49000 下 Softmax 最佳区域在 learning rate 约 5e-7 到 7.5e-7、reg 约 5e3 附近。正则过强会压低训练与验证准确率。

### 5.4 TwoLayerNet 配置对比

![TwoLayerNet 配置对比](figures/two_layer_config_comparison.png)

TwoLayerNet 对超参数更敏感。较大的 hidden_dim 并不必然更好，训练 epoch、正则强度与学习率需要共同匹配。

### 5.5 数据规模影响

![数据规模影响](figures/data_scale_effect.png)

更多数据对 Softmax 和 TwoLayerNet 更有帮助；kNN 的瓶颈主要是原始像素距离质量，因此提升有限。

### 5.6 运行耗时

![运行耗时对比](figures/runtime_comparison.png)

kNN 几乎没有训练成本，但预测时要与训练集逐一比较，规模变大后推理成本线性上升。TwoLayerNet 的主要成本在训练阶段。

### 5.7 最佳 TwoLayerNet 训练损失与混淆矩阵

![TwoLayerNet 损失曲线](figures/two_layer_loss_curve.png)

![混淆矩阵](figures/confusion_matrix_best_model.png)

本次用于诊断图的重新训练配置为：{diag['config']}，测试准确率为 {diag['test_acc']:.3f}。

## 6. 为什么 TwoLayerNet 更好

kNN 依赖原始像素 L2 距离，但图像语义与像素距离并不稳定：光照、背景、位置变化都会显著改变像素距离。Softmax 学习每个类别的线性模板，能比 kNN 更稳定，但决策边界仍是线性的。TwoLayerNet 通过隐藏层和 ReLU 组合多个线性变换，可以表达非线性决策边界，因此在 CIFAR-10 原始像素上显著领先。

## 7. 当前实验局限

- 仍然使用 raw pixels，没有使用 HOG / color histogram 等特征。
- TwoLayerNet 只做了代表性配置搜索，还不是极致调参。
- 测试集只取 1000 张用于快速复现，完整 CIFAR-10 测试集为 10000 张。
- 没有训练 CNN，因此还没有利用图像局部空间结构。

## 8. 下一步计划

1. 跑 `features.ipynb`：加入 HOG 和颜色直方图特征。
2. 跑 `FullyConnectedNets.ipynb`：比较多层网络、Dropout、Batch Normalization。
3. 增加 TwoLayerNet epoch 和更细网格搜索。
4. 后续进入 CNN，在 CIFAR-10 上进一步提升准确率。
"""
    (RESULTS / "week5_full_reproduction_report.md").write_text(md, encoding="utf-8")


def add_textbox(slide, text, x, y, w, h, size=20, bold=False, color=(42, 46, 52), align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    if align:
        p.alignment = align
    return box


def add_title(slide, kicker, title):
    add_textbox(slide, kicker.upper(), 0.55, 0.28, 2.5, 0.26, size=9, bold=True, color=(95, 111, 130))
    add_textbox(slide, title, 0.55, 0.55, 11.4, 0.55, size=24, bold=True, color=(30, 35, 40))


def add_bullets(slide, bullets, x=0.75, y=1.35, w=5.5, h=4.8, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.space_after = Pt(8)
    return box


def add_image(slide, filename, x, y, w, h=None):
    p = FIGURES / filename
    if h is None:
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def make_ppt(all_results, df):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide_bg(slide, color=(248, 247, 244)):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    best = best_overall(df)
    n49 = all_results["N49000"]
    slides = []

    s = prs.slides.add_slide(blank); slide_bg(s); slides.append(s)
    add_textbox(s, "Week5 CS231n Assignment 1", 0.7, 1.0, 10.5, 0.6, 34, True)
    add_textbox(s, "CIFAR-10 图像分类本地复现实验：kNN、Softmax、TwoLayerNet", 0.72, 1.75, 10.2, 0.4, 18)
    add_textbox(s, f"最佳测试准确率：{best['model']} / {best['scale']} / {best['test_acc']:.3f}", 0.72, 2.55, 8.5, 0.45, 20, True, (79, 137, 104))
    add_textbox(s, "数据规模：10k / 20k / 49k 训练样本；验证与测试各 1k", 0.72, 3.05, 8.5, 0.35, 15)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "background", "任务背景：从图像分类 pipeline 进入神经网络")
    add_bullets(s, ["CS231n A1 关注图像分类的完整工程链路", "从 kNN 到 Softmax，再到两层神经网络", "核心学习点：数据划分、loss、gradient、regularization、hyperparameter tuning", "本实验不是提交作业，而是本地复现与组会汇报"], w=11.5)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "dataset", "CIFAR-10：10 类 32×32 彩色图像")
    add_bullets(s, ["原始数据：50,000 train + 10,000 test", "官方 A1 常用划分：49,000 train + 1,000 val", "本地复现实验对比：N10000、N20000、N49000", "快速评估阶段使用 1,000 张 test 以控制运行时间"], x=0.75, y=1.25, w=5.2)
    add_image(s, "data_scale_effect.png", 6.3, 1.25, 6.2, 4.1)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "workflow", "实验流程：训练、调参、验证、测试与可视化")
    add_bullets(s, ["1. 加载 CIFAR-10 并做 train/val/test split", "2. 训练 kNN、Softmax、TwoLayerNet", "3. 在验证集上选择超参数", "4. 用测试集报告泛化性能", "5. 汇总 JSON、Markdown 报告、PNG 图表和 PPT"], w=11.5)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "knn", "kNN：简单但被原始像素距离限制")
    add_bullets(s, [f"N49000 最佳 k={n49['knn']['best_k']}", f"验证准确率 {n49['knn']['best_val_acc']:.3f}，测试准确率 {n49['knn']['test_acc']:.3f}", "几乎没有训练成本，但推理要比较所有训练样本", "在高维原始像素空间中，L2 距离不稳定"], x=0.7, y=1.25, w=5.0)
    add_image(s, "runtime_comparison.png", 6.1, 1.3, 6.5, 4.4)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "softmax", "Softmax：线性分类器受益于调参和更多数据")
    add_bullets(s, [f"N49000 最佳 lr={n49['softmax']['best']['learning_rate']:.1e}", f"最佳 reg={n49['softmax']['best']['reg']:.0f}", f"验证准确率 {n49['softmax']['best']['val_acc']:.3f}，测试准确率 {n49['softmax']['best']['test_acc']:.3f}", "线性决策边界仍限制了上限"], x=0.65, y=1.25, w=4.5)
    add_image(s, "softmax_hyperparam_heatmap.png", 5.35, 1.15, 7.2, 4.8)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "two-layer net", "TwoLayerNet：非线性隐藏层带来明显提升")
    add_bullets(s, [f"N49000 最佳验证配置：H={n49['two_layer_net']['best']['hidden_dim']}, lr={n49['two_layer_net']['best']['learning_rate']:.1e}, reg={n49['two_layer_net']['best']['reg']}", f"验证准确率 {n49['two_layer_net']['best']['best_val_acc']:.3f}", "隐藏层 + ReLU 学习非线性特征组合", "模型容量越大不一定越好，需要正则化与学习率匹配"], x=0.55, y=1.18, w=4.9, size=14)
    add_image(s, "two_layer_config_comparison.png", 5.15, 1.1, 7.7, 4.9)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "comparison", "模型对比：TwoLayerNet 在 raw pixels 上领先")
    add_image(s, "model_test_accuracy_comparison.png", 0.9, 1.25, 11.6, 5.0)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "scale effect", "数据规模影响：线性与非线性模型更能利用更多数据")
    add_image(s, "data_scale_effect.png", 0.75, 1.2, 7.1, 4.7)
    add_bullets(s, ["kNN：提升有限，瓶颈在像素距离", "Softmax：更多数据带来稳定小幅提升", "TwoLayerNet：49k 后突破 0.5 附近", "测试集波动不替代验证集选模"], x=8.2, y=1.45, w=4.2, size=15)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "diagnostics", "诊断图：loss 下降与类别混淆")
    add_image(s, "two_layer_loss_curve.png", 0.75, 1.15, 5.6, 3.25)
    add_image(s, "confusion_matrix_best_model.png", 6.7, 1.0, 5.5, 4.7)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "best model", "最佳模型：TwoLayerNet 是本轮复现的主结果")
    add_bullets(s, [f"按 validation：N49000 TwoLayerNet cfg4，val={n49['two_layer_net']['best']['best_val_acc']:.3f}, test={n49['two_layer_net']['best']['test_acc']:.3f}", "按 test：部分配置测试准确率可到 0.510，但应以验证集选模", f"Softmax 最佳 test={n49['softmax']['best']['test_acc']:.3f}", f"kNN 最佳 test={n49['knn']['test_acc']:.3f}"], w=11.4)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "takeaways", "关键学习收获")
    add_bullets(s, ["图像分类 pipeline 的核心是数据划分 + 训练 + 验证调参 + 测试评估", "向量化 NumPy 实现决定了实验能否高效运行", "正则化和学习率会直接决定泛化性能", "raw pixels 能跑通流程，但不是好的视觉表示", "非线性模型明显优于线性模型，但还不是 CNN"], w=11.4)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "limits", "当前局限：还没有进入真正强的图像表示")
    add_bullets(s, ["尚未充分跑 HOG / color histogram features", "FullyConnectedNet 的 dropout / batch normalization 还未展开", "测试集使用 1,000 张快速评估，后续可扩到完整 10,000 张", "未使用 CNN，因此没有利用图像局部空间结构"], w=11.4)

    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "next", "下一步计划：从 A1 走向更强模型")
    add_bullets(s, ["1. 跑 features.ipynb：HOG + color histogram", "2. 跑 FullyConnectedNets.ipynb：多层网络、Dropout、BatchNorm", "3. 扩大 TwoLayerNet epoch 与更细调参", "4. 进入 CNN：利用局部感受野和权重共享"], w=11.4)

    out = RESULTS / "week5_cs231n_assignment1_report.pptx"
    prs.save(out)

    outline = """# Week5 CS231n Assignment 1 PPT 大纲

1. 标题页：CIFAR-10 图像分类本地复现实验
2. 任务背景：CS231n A1 的图像分类 pipeline
3. 数据集介绍：CIFAR-10 与 10k/20k/49k 规模
4. 实验流程：数据、训练、调参、测试、可视化
5. kNN：最近邻思想、结果与推理成本
6. Softmax：线性分类器与超参数热力图
7. TwoLayerNet：非线性隐藏层与配置对比
8. 模型整体测试准确率对比
9. 数据规模影响
10. 训练损失与混淆矩阵诊断
11. 最佳模型总结
12. 关键学习收获
13. 当前局限
14. 下一步计划
"""
    (RESULTS / "week5_cs231n_assignment1_ppt_outline.md").write_text(outline, encoding="utf-8")
    return out


def main():
    all_results = load_all_results()
    write_json(RESULTS / "week5_49000_reproduction_summary.json", all_results["N49000"])
    write_json(RESULTS / "week5_all_scales_summary.json", all_results)
    df, diag = make_figures(all_results)
    df.to_csv(RESULTS / "week5_model_comparison_table.csv", index=False, encoding="utf-8-sig")
    write_report(all_results, df, diag)
    pptx = make_ppt(all_results, df)
    print("Generated:")
    print(RESULTS / "week5_49000_reproduction_summary.json")
    print(RESULTS / "week5_all_scales_summary.json")
    print(RESULTS / "week5_full_reproduction_report.md")
    print(pptx)


if __name__ == "__main__":
    main()
